from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs, title, subtitle_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    subtitle = slide.placeholders[1].text_frame
    subtitle.clear()

    for i, line in enumerate(subtitle_lines):
        p = subtitle.paragraphs[0] if i == 0 else subtitle.add_paragraph()
        p.text = line
        p.font.size = Pt(20 if i == 0 else 16)
        p.alignment = PP_ALIGN.CENTER


def add_bullet_slide(prs, title, bullets):
    """
    bullets: list có thể chứa:
      - string (bullet cấp 1)
      - tuple(level, text) với level: 0,1,2...
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item

        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(22 if level == 0 else 18)


def build_presentation(output_path="Bao_cao_Tuan10_Bao_mat_ung_dung.pptx"):
    prs = Presentation()

    # Slide 1
    add_title_slide(
        prs,
        "BÁO CÁO HỌC TẬP TUẦN 10: BẢO MẬT ỨNG DỤNG",
        [
            "Security in AI-Augmented Development",
            "Khóa học AI-Augmented Project-Based Programming - Trường DNTU",
            "Ngày báo cáo: 07/05/2026",
        ],
    )

    # Slide 2
    add_bullet_slide(prs, "Tổng quan & Mục tiêu", [
        "Authentication (Xác thực) & Authorization (Phân quyền).",
        "Data Privacy (Quyền riêng tư dữ liệu).",
        "Tiêu chuẩn bảo mật API (OWASP Top 10).",
        "Ứng dụng AI Tools trong bảo mật.",
    ])

    # Slide 3
    add_bullet_slide(prs, "Authentication - Xác thực người dùng", [
        "Định nghĩa: Quá trình xác nhận danh tính người dùng (Lớp bảo vệ đầu tiên).",
        "Các phương pháp phổ biến:",
        (1, "JWT (JSON Web Token): Phổ biến cho REST API."),
        (1, "OAuth 2.0 / OpenID Connect: Đăng nhập qua bên thứ ba (Google, GitHub)."),
        (1, "MFA & Biometric: Xác thực đa nhân tố và sinh trắc học."),
    ])

    # Slide 4
    add_bullet_slide(prs, "Tìm hiểu sâu về JWT (JSON Web Token)", [
        "Cấu trúc 3 phần: Header, Payload, Signature.",
        "Cơ chế Token:",
        (1, "Access Token: Thời hạn ngắn (15p - 1h)."),
        (1, "Refresh Token: Thời hạn dài (7 - 30 ngày)."),
        "Best Practice: Lưu trong httpOnly cookie, tránh dùng localStorage.",
    ])

    # Slide 5
    add_bullet_slide(prs, "Authorization - Phân quyền truy cập", [
        "Khái niệm: Xác định quyền hạn sau khi đã xác thực.",
        "Các mô hình chính:",
        (1, "RBAC: Phân quyền theo vai trò (Admin, User, Guest)."),
        (1, "ABAC: Phân quyền theo thuộc tính chi tiết."),
        (1, "ACL: Danh sách kiểm soát truy cập cụ thể."),
    ])

    # Slide 6
    add_bullet_slide(prs, "Triển khai trong NestJS", [
        "Authentication: Sử dụng @nestjs/jwt và @nestjs/passport.",
        "Authorization (RBAC):",
        (1, "Dùng @Roles() decorator."),
        (1, "Triển khai RolesGuard để kiểm tra quyền trước request."),
        "AI Support: Dùng Claude/Copilot để tạo boilerplate nhanh.",
    ])

    # Slide 7
    add_bullet_slide(prs, "Data Privacy - Bảo vệ dữ liệu cá nhân", [
        "Nguyên tắc cốt lõi:",
        (1, "Data Minimization: Chỉ thu thập dữ liệu cần thiết."),
        (1, "Encryption: Mã hóa dữ liệu tĩnh (At Rest) và dữ liệu đang truyền (In Transit - HTTPS)."),
        (1, "Right to Erasure: Quyền yêu cầu xóa dữ liệu (GDPR)."),
        (1, "Audit Logging: Ghi nhật ký truy cập dữ liệu nhạy cảm."),
    ])

    # Slide 8
    add_bullet_slide(prs, "OWASP Top 10 & Bài tập thực hành", [
        "Lỗ hổng cần lưu ý: Broken Access Control, Cryptographic Failures, Identification Failures.",
        "Yêu cầu thực hành:",
        (1, "Hash password với bcrypt (salt rounds >= 12)."),
        (1, "Áp dụng Rate Limiting cho login."),
        (1, "Cấu hình CORS, Helmet và security headers."),
        (1, "Unit test cho Auth logic (Coverage >= 80%)."),
    ])

    # Slide 9
    add_bullet_slide(prs, "Kết luận & Hướng phát triển", [
        "Kết nối: Liên kết với NestJS (Tuần 4), Database (Tuần 6) và Testing (Tuần 8).",
        "Tương lai:",
        (1, "Tìm hiểu Zero Trust Architecture & DevSecOps."),
        (1, "Thực hành quét lỗ hổng với OWASP ZAP, Snyk."),
    ])

    prs.save(output_path)
    print(f"Da tao file: {output_path}")


if __name__ == "__main__":
    build_presentation()
